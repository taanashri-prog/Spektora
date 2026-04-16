import os
import io
from flask import Flask, render_template_string, request, send_file
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

app = Flask(__name__)

# --- UI & UX DESIGN (HTML/CSS/JS) ---
HTML_TEMPLATE = """
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <title>Spektora | Your Voice. Our Vision.</title>
    <style>
        :root {
            --bg: #0D0D0D;
            --accent: #D4AF37;
            --card: #1A1A1A;
            --text: #FFFFFF;
        }
        .light-theme { --bg: #F5F5F5; --card: #FFFFFF; --text: #1A1A1A; }

        body { font-family: 'Arial', sans-serif; background: var(--bg); color: var(--text); margin: 0; overflow: hidden; transition: 0.5s; }

        /* SPLASH SCREEN & LOGO ANIMATION */
        #splash {
            position: fixed; width: 100%; height: 100vh; background: #000;
            display: flex; flex-direction: column; justify-content: center; align-items: center;
            z-index: 9999; transition: 1s cubic-bezier(0.7, 0, 0.3, 1);
        }

        #logo-container { position: relative; }
        .logo-main {
            font-size: 80px; font-weight: bold; color: var(--accent);
            transition: 1.2s cubic-bezier(0.7, 0, 0.3, 1);
            text-shadow: 0 0 20px rgba(212, 175, 55, 0.6);
        }

        /* Glitter particles around the logo */
        .sparkle {
            position: absolute; width: 4px; height: 4px; background: var(--accent);
            border-radius: 50%; pointer-events: none; opacity: 0;
            animation: glitter-move 2s infinite ease-in-out;
        }
        @keyframes glitter-move {
            0% { transform: translate(0,0); opacity: 0; }
            50% { opacity: 1; }
            100% { transform: translate(calc(Math.random()*100px - 50px), -100px); opacity: 0; }
        }

        /* MAIN INTERFACE */
        #main-app { display: flex; height: 100vh; opacity: 0; transition: 1s; }

        .sidebar { width: 320px; background: var(--card); padding: 40px; border-right: 1px solid rgba(212,175,55,0.1); position: relative; }
        .spark-box {
            background: rgba(212,175,55,0.05); border: 1px solid var(--accent);
            padding: 20px; border-radius: 20px; margin-top: 60px; position: relative;
        }
        .spark-title { font-weight: bold; color: var(--accent); letter-spacing: 1px; margin-bottom: 5px; }

        .potted-plant { position: absolute; bottom: 30px; left: 30px; font-size: 80px; opacity: 0.9; }

        .engine-room { flex-grow: 1; display: flex; flex-direction: column; justify-content: center; align-items: center; padding: 40px; }
        .glass-card {
            background: var(--card); padding: 50px; border-radius: 30px;
            box-shadow: 0 25px 50px rgba(0,0,0,0.5); width: 100%; max-width: 600px;
            border: 1px solid rgba(255,255,255,0.05);
        }

        input, textarea, select {
            width: 100%; padding: 18px; margin: 12px 0; border-radius: 12px;
            background: rgba(0,0,0,0.2); color: white; border: 1px solid #333; box-sizing: border-box;
        }

        .btn-gen {
            background: var(--accent); color: black; border: none; padding: 22px;
            border-radius: 15px; font-weight: 800; width: 100%; cursor: pointer;
            font-size: 18px; letter-spacing: 2px; transition: 0.3s ease;
        }
    </style>
</head>
<body>
    <div id="splash">
        <div id="logo-container">
            <div class="logo-main">Spektora</div>
        </div>
    </div>

    <div id="main-app">
        <div class="sidebar">
            <div class="logo-main" style="font-size: 30px; margin-bottom: 40px;">Spektora</div>
            <div class="spark-box">
                <div class="spark-title">Presentation Style</div>
                <select id="style-select">
                    <option value="modern">Modern</option>
                    <option value="corporate">Corporate</option>
                    <option value="creative">Creative</option>
                </select>
            </div>

            <div class="spark-box">
                <div class="spark-title">Content Focus</div>
                <select id="focus-select">
                    <option value="executive">Executive Summary</option>
                    <option value="technical">Technical Deep Dive</option>
                    <option value="marketing">Marketing Pitch</option>
                </select>
            </div>

            <div class="spark-box">
                <div class="spark-title">Theme</div>
                <label class="switch">
                    <input type="checkbox" id="theme-toggle">
                    <span class="slider round"></span>
                </label>
                <span id="theme-text" style="margin-left: 10px;">Dark</span>
            </div>
            <div class="potted-plant">🌱</div>
        </div>

        <div class="engine-room">
            <div class="glass-card">
                <h1 style="color: var(--accent); text-align: center; margin-bottom: 30px;">Generate Your Presentation</h1>
                <form action="/generate_pptx" method="post">
                    <label for="topic" style="color: var(--text);">Topic:</label>
                    <input type="text" id="topic" name="topic" placeholder="E.g., Future of AI" required>

                    <label for="audience" style="color: var(--text);">Audience:</label>
                    <input type="text" id="audience" name="audience" placeholder="E.g., Investors, Technical Team" required>

                    <label for="key_points" style="color: var(--text);">Key Points (comma-separated):</label>
                    <textarea id="key_points" name="key_points" rows="4" placeholder="E.g., AI in healthcare, ethical considerations, market potential" required></textarea>

                    <button type="submit" class="btn-gen">GENERATE PRESENTATION</button>
                </form>
            </div>
        </div>
    </div>

    <script>
        // SPLASH SCREEN LOGIC
        window.addEventListener('load', () => {
            setTimeout(() => {
                document.getElementById('splash').style.opacity = '0';
                document.getElementById('splash').style.transform = 'scale(1.1)';
                setTimeout(() => {
                    document.getElementById('splash').style.display = 'none';
                    document.getElementById('main-app').style.opacity = '1';
                }, 1000); // Wait for fade-out to complete
            }, 2000); // Show splash for 2 seconds

            // Generate some sparkles
            const logoContainer = document.getElementById('logo-container');
            for (let i = 0; i < 15; i++) {
                let sparkle = document.createElement('div');
                sparkle.classList.add('sparkle');
                sparkle.style.top = `${Math.random() * 100}%`;
                sparkle.style.left = `${Math.random() * 100}%`;
                sparkle.style.animationDelay = `${Math.random() * 2}s`;
                logoContainer.appendChild(sparkle);
            }
        });

        // THEME TOGGLE LOGIC
        const themeToggle = document.getElementById('theme-toggle');
        const themeText = document.getElementById('theme-text');

        themeToggle.addEventListener('change', () => {
            document.body.classList.toggle('light-theme');
            themeText.textContent = themeToggle.checked ? 'Light' : 'Dark';
        });

        // Store selected options in localStorage for persistence (optional)
        document.getElementById('style-select').addEventListener('change', function() {
            localStorage.setItem('presentationStyle', this.value);
        });

        document.getElementById('focus-select').addEventListener('change', function() {
            localStorage.setItem('contentFocus', this.value);
        });

        // Load saved options
        window.addEventListener('DOMContentLoaded', (event) => {
            const savedStyle = localStorage.getItem('presentationStyle');
            if (savedStyle) {
                document.getElementById('style-select').value = savedStyle;
            }
            const savedFocus = localStorage.getItem('contentFocus');
            if (savedFocus) {
                document.getElementById('focus-select').value = savedFocus;
            }
        });
    </script>
</body>
</html>
"""

@app.route('/')
def index():
    return render_template_string(HTML_TEMPLATE)

@app.route('/generate_pptx', methods=['POST'])
def generate_pptx():
    topic = request.form['topic']
    audience = request.form['audience']
    key_points_str = request.form['key_points']
    key_points = [kp.strip() for kp in key_points_str.split(',') if kp.strip()]

    prs = Presentation()

    # Title Slide
    title_slide_layout = prs.slide_layouts[0] # Usually the title slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = topic
    subtitle.text = f"A presentation for {audience}"

    # Set title font size and color
    title.text_frame.paragraphs[0].font.size = Pt(48)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xD4, 0xAF, 0x37) # Gold color

    # Add a content slide for each key point
    bullet_slide_layout = prs.slide_layouts[1] # Title and Content layout

    for i, point in enumerate(key_points):
        slide = prs.slides.add_slide(bullet_slide_layout)
        title = slide.shapes.title
        body = slide.placeholders[1]

        title.text = f"Key Point {i+1}: {point}"

        tf = body.text_frame
        p = tf.add_paragraph()
        p.text = "[Elaborate on this point here. Add details, examples, data.]"
        p.level = 0

        p2 = tf.add_paragraph()
        p2.text = "- Sub-point 1"
        p2.level = 1

        p3 = tf.add_paragraph()
        p3.text = "- Sub-point 2"
        p3.level = 1

    # Example of adding a chart slide (requires actual data generation logic)
    # This is a placeholder and would need proper data integration
    chart_slide_layout = prs.slide_layouts[5] # Title Only layout (can be adapted for charts)
    slide = prs.slides.add_slide(chart_slide_layout)
    title = slide.shapes.title
    title.text = "Example Data Visualization"

    # Create chart data
    chart_data = CategoryChartData()
    chart_data.categories = ['East', 'West', 'Midwest']
    chart_data.add_series('Series 1', (19.2, 21.4, 16.7))

    x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    )

    # Save presentation to a byte stream
    pptx_file = io.BytesIO()
    prs.save(pptx_file)
    pptx_file.seek(0)

    return send_file(pptx_file, as_attachment=True, download_name=f"{topic.replace(' ', '_')}_Presentation.pptx", mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')

if __name__ == '__main__':
    # You can change the port if 5000 is in use
    app.run(host='0.0.0.0', port=5000)
